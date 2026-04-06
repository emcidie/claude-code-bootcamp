from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x1a, 0x1a, 0x2e)
DARK_BLUE = RGBColor(0x16, 0x21, 0x3e)
ORANGE = RGBColor(0xe0, 0x7a, 0x3a)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x9c, 0xa3, 0xaf)
LIGHT_BG = RGBColor(0xf3, 0xf4, 0xf6)
GREEN = RGBColor(0x10, 0xb9, 0x81)
RED = RGBColor(0xdc, 0x26, 0x26)

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, left, top, width, height, items, font_size=20, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(12)
        p.level = 0
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, title_color=ORANGE, body_color=WHITE):
    rect = add_rounded_rect(slide, left, top, width, height, DARK_BLUE)
    add_text_box(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, title, font_size=18, bold=True, color=title_color)
    add_text_box(slide, left + 0.3, top + 0.7, width - 0.6, height - 0.9, body, font_size=14, color=body_color)

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

add_text_box(slide, 1, 0.8, 11, 0.6, 'JAPAN CMT', font_size=24, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.5, 11, 1.2, 'Introduction to Claude', font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 0.6, 'AI That Works With You', font_size=28, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Divider line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(2.3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 1, 4.2, 11, 0.5, 'Claude Code Bootcamp 2026', font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.8, 11, 0.5, '20-Minute Overview', font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: What is AI?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'What is AI?', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 0.8, 1.6, 11, 0.8, 'AI is software that learned from billions of texts — books, websites, code —\nso it can understand your questions and generate helpful responses.', font_size=22, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.8, 3.6, 2.2, 'Think of it like...', 'A colleague who has read every programming textbook, every Stack Overflow answer, and every tutorial ever written.')
add_card(slide, 4.7, 2.8, 3.6, 2.2, 'What it can do', 'Write code, explain concepts, debug errors, create files, answer questions — all in plain English.')
add_card(slide, 8.6, 2.8, 3.6, 2.2, 'What it cannot do', 'It doesn\'t "think" like humans. It predicts useful responses based on patterns. It can make mistakes.')

add_text_box(slide, 0.8, 5.4, 11, 0.6, 'Key takeaway: AI is a powerful tool, not magic. The better you communicate with it, the better results you get.', font_size=16, bold=True, color=ORANGE)

# ============================================================
# SLIDE 3: Meet Claude
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'Meet Claude', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 0.8, 1.6, 6, 0.6, 'Claude is an AI assistant made by Anthropic.', font_size=22, color=LIGHT_GRAY)

items = [
    '  Helpful — designed to give accurate, useful answers',
    '  Harmless — built with safety guardrails',
    '  Honest — tells you when it\'s unsure or doesn\'t know',
    '  Versatile — writes code, explains concepts, creates content',
    '  Conversational — remembers context within a session'
]
add_bullet_slide(slide, 0.8, 2.3, 7, 3.5, items, font_size=22, color=WHITE)

# Right side: model comparison
add_rounded_rect(slide, 8.2, 1.6, 4.3, 4.5, DARK_BLUE)
add_text_box(slide, 8.5, 1.8, 3.8, 0.5, 'Claude Model Family', font_size=18, bold=True, color=ORANGE)

models = [
    ('Haiku', 'Fast & affordable\nGreat for quick tasks', GREEN),
    ('Sonnet', 'Balanced performance\nIdeal for most work', ORANGE),
    ('Opus', 'Most capable\nComplex reasoning', RGBColor(0xa7, 0x8b, 0xfa)),
]
for i, (name, desc, color) in enumerate(models):
    y = 2.5 + i * 1.15
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(y + 0.05), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, 8.9, y - 0.05, 1.5, 0.4, name, font_size=18, bold=True, color=WHITE)
    add_text_box(slide, 8.9, y + 0.3, 3.2, 0.6, desc, font_size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 4: Claude vs Claude Code
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'Claude vs Claude Code', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

# Left card - Claude
add_rounded_rect(slide, 0.8, 1.8, 5.5, 4.2, DARK_BLUE)
add_text_box(slide, 1.2, 2.0, 4.8, 0.5, 'Claude (Chatbot)', font_size=26, bold=True, color=ORANGE)
claude_items = [
    '  Chat on claude.ai or the app',
    '  Ask questions, get answers',
    '  Copy & paste code manually',
    '  Great for learning & research',
    '  Like texting a smart friend'
]
add_bullet_slide(slide, 1.2, 2.7, 4.8, 3.0, claude_items, font_size=18, color=WHITE)

# Right card - Claude Code
add_rounded_rect(slide, 7.0, 1.8, 5.5, 4.2, DARK_BLUE)
add_text_box(slide, 7.4, 2.0, 4.8, 0.5, 'Claude Code (Dev Tool)', font_size=26, bold=True, color=GREEN)
cc_items = [
    '  Works in your terminal',
    '  Reads & edits your files directly',
    '  Runs commands for you',
    '  Builds entire applications',
    '  Like a developer sitting next to you'
]
add_bullet_slide(slide, 7.4, 2.7, 4.8, 3.0, cc_items, font_size=18, color=WHITE)

# Arrow
add_text_box(slide, 5.8, 3.2, 1.5, 1, '>>>', font_size=40, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.2, 11, 0.5, 'Today we\'re learning Claude Code — the hands-on developer tool!', font_size=18, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: What Can You Build?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'What Can You Build with Claude Code?', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 0.8, 1.5, 11, 0.6, 'Even as a complete beginner, you can build real things today:', font_size=20, color=LIGHT_GRAY)

grid = [
    ('Websites', 'Personal portfolios,\nlanding pages, blogs'),
    ('Web Apps', 'Dashboards, tools,\ninteractive sites'),
    ('APIs', 'Backend services\nthat power apps'),
    ('Automations', 'Scripts that save\nyou hours of work'),
    ('Games', 'Browser-based games\nwith HTML5 & JS'),
    ('AI Apps', 'Chatbots, tutors,\nAI-powered tools'),
]

for i, (title, desc) in enumerate(grid):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 2.3 + row * 2.0
    add_rounded_rect(slide, x, y, 3.6, 1.7, DARK_BLUE)
    add_text_box(slide, x + 0.3, y + 0.2, 3.0, 0.4, title, font_size=22, bold=True, color=ORANGE)
    add_text_box(slide, x + 0.3, y + 0.7, 3.0, 0.8, desc, font_size=15, color=WHITE)

# ============================================================
# SLIDE 6: How Claude Code Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'How Claude Code Works', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 0.8, 1.5, 11, 0.6, 'The workflow is simple: You describe what you want. Claude Code does it.', font_size=20, color=LIGHT_GRAY)

steps = [
    ('1', 'You Describe', 'Tell Claude Code what\nyou want in plain English', ORANGE),
    ('2', 'Claude Writes', 'It generates code,\ncreates or edits files', RGBColor(0xa7, 0x8b, 0xfa)),
    ('3', 'You Review', 'Check the changes\nand approve them', RGBColor(0x60, 0xa5, 0xfa)),
    ('4', 'You Refine', 'Ask for tweaks until\nit\'s exactly right', GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.8 + i * 3.1
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(2.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + 0.9, 2.35, 0.7, 0.6, num, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x, 3.2, 2.8, 0.5, title, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, 3.7, 2.8, 0.8, desc, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

# Example box
add_rounded_rect(slide, 1.5, 4.8, 10.3, 1.5, DARK_BLUE)
add_text_box(slide, 1.8, 4.9, 3, 0.4, 'Example:', font_size=16, bold=True, color=ORANGE)
add_text_box(slide, 1.8, 5.3, 9.5, 0.9, 'You say: "Create a personal portfolio webpage with my name, an about section, and a contact form"\nClaude Code: Creates the complete HTML, CSS, and JavaScript file — ready to open in a browser!', font_size=16, color=WHITE)

# ============================================================
# SLIDE 7: Key Concepts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'Key Concepts to Remember', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

concepts = [
    ('Prompts', 'The messages you send to Claude Code.\nBetter prompts = better results.', 'Be specific: "Create an HTML page with a blue header" beats "make a website"'),
    ('Context', 'Claude Code sees the files in your folder.\nIt uses that context to give better answers.', 'Always work in the right project folder'),
    ('Permissions', 'Claude Code asks before changing your files.\nYou stay in control at all times.', 'Review changes before approving them'),
    ('Slash Commands', 'Built-in shortcuts like /help, /clear, /compact.\nThey help you manage your session.', 'Use /compact when conversations get long'),
]

for i, (title, desc, tip) in enumerate(concepts):
    y = 1.6 + i * 1.35
    add_rounded_rect(slide, 0.8, y, 11.7, 1.15, DARK_BLUE)
    add_text_box(slide, 1.1, y + 0.1, 2.2, 0.4, title, font_size=20, bold=True, color=ORANGE)
    add_text_box(slide, 3.5, y + 0.1, 4.5, 0.9, desc, font_size=15, color=WHITE)
    add_text_box(slide, 8.5, y + 0.15, 3.8, 0.85, tip, font_size=14, color=GREEN)

# ============================================================
# SLIDE 8: Bootcamp Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.5, 11, 0.8, 'Today\'s Bootcamp Agenda', font_size=40, bold=True, color=WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 0.8, 1.5, 11, 0.5, '2.5 hours of hands-on learning with an AI tutor by your side', font_size=20, color=LIGHT_GRAY)

modules = [
    ('1', 'What is AI & Claude Code?', '15 min'),
    ('2', 'Setting Up Claude Code', '15 min'),
    ('3', 'Your First Conversation', '15 min'),
    ('4', 'Prompts & Skills', '20 min'),
    ('5', 'File Operations', '20 min'),
    ('6', 'Building a Web Page', '25 min'),
    ('7', 'Git & GitHub', '20 min'),
    ('8', 'Tips & Next Steps', '20 min'),
]

for i, (num, title, duration) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.2 + row * 1.15

    add_rounded_rect(slide, x, y, 5.8, 0.95, DARK_BLUE)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.2), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    add_text_box(slide, x + 0.15, y + 0.22, 0.55, 0.5, num, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + 0.9, y + 0.2, 3.5, 0.5, title, font_size=18, bold=False, color=WHITE)
    add_text_box(slide, x + 4.5, y + 0.25, 1.1, 0.4, duration, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 9: Let's Get Started
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1, 1.5, 11, 1.2, 'Let\'s Get Started!', font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.0), Inches(2.3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_text_box(slide, 1, 3.5, 11, 0.6, 'Open your browser and go to:', font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.8, 'emcidie.com/bootcamp/', font_size=36, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.2, 11, 0.6, 'Enter your name and start Module 1', font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.0, 11, 0.5, 'Your AI tutor is ready to help you every step of the way', font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = '/Users/mikeduff/claude-code-bootcamp/Japan_CMT_Introduction_to_Claude.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
