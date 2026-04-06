from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Accenture-style colors
PURPLE = RGBColor(0xA1, 0x00, 0xFF)
DARK_PURPLE = RGBColor(0x7B, 0x00, 0xC4)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xB4, 0x55)
BLUE = RGBColor(0x00, 0x70, 0xC0)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=18, color=DARK_GRAY, bullet_color=PURPLE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(10)
    return tf

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_gt_symbol(slide, x, y, size=60, color=PURPLE):
    """Add Accenture-style > symbol"""
    add_text_box(slide, x, y, 1.5, 1.5, '>', font_size=size, bold=True, color=color)

def add_purple_bar(slide):
    """Add the purple bottom bar - Accenture signature"""
    add_rect(slide, 0, 7.1, 13.333, 0.4, PURPLE)

def add_top_bar(slide):
    """Add thin purple top accent line"""
    add_rect(slide, 0, 0, 13.333, 0.06, PURPLE)

def add_card_white(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

# ============================================================
# SLIDE 1: Title Slide (Purple background, Accenture style)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PURPLE)

# Large > symbol as background element
add_text_box(slide, 9, -0.5, 5, 8, '>', font_size=350, bold=True, color=DARK_PURPLE, alignment=PP_ALIGN.RIGHT)

add_text_box(slide, 1, 1.0, 8, 0.5, 'JAPAN CMT', font_size=20, bold=True, color=WHITE)

# Thin line
add_rect(slide, 1, 1.6, 1.5, 0.04, WHITE)

add_text_box(slide, 1, 2.0, 10, 1.5, 'Introduction\nto Claude', font_size=56, bold=True, color=WHITE)
add_text_box(slide, 1, 4.2, 8, 0.6, 'AI That Works With You', font_size=26, bold=False, color=WHITE)

add_text_box(slide, 1, 5.5, 6, 0.4, 'Claude Code Bootcamp 2026', font_size=16, color=WHITE)
add_text_box(slide, 1, 5.9, 6, 0.4, '20-Minute Overview', font_size=14, color=WHITE)

# ============================================================
# SLIDE 2: What is AI?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, 'What is AI?', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)

add_text_box(slide, 0.8, 1.4, 11, 0.6, 'AI is software that learned from billions of texts — books, websites, code —\nso it can understand your questions and generate helpful responses.', font_size=18, color=MID_GRAY)

# Three cards on light gray background
add_rect(slide, 0.5, 2.3, 12.3, 4.5, LIGHT_GRAY)

# Card 1
add_card_white(slide, 0.8, 2.6, 3.6, 3.8)
add_gt_symbol(slide, 0.9, 2.6, 40, PURPLE)
add_text_box(slide, 0.9, 3.3, 3.2, 0.4, 'Think of it like...', font_size=20, bold=True, color=BLACK)
add_text_box(slide, 0.9, 3.8, 3.2, 1.5, 'A colleague who has read every programming textbook, every Stack Overflow answer, and every tutorial ever written — and can recall it instantly.', font_size=15, color=MID_GRAY)

# Card 2
add_card_white(slide, 4.8, 2.6, 3.6, 3.8)
add_gt_symbol(slide, 4.9, 2.6, 40, PURPLE)
add_text_box(slide, 4.9, 3.3, 3.2, 0.4, 'What it can do', font_size=20, bold=True, color=BLACK)
add_text_box(slide, 4.9, 3.8, 3.2, 1.5, 'Write code, explain concepts, debug errors, create and edit files, answer questions — all from plain English instructions.', font_size=15, color=MID_GRAY)

# Card 3
add_card_white(slide, 8.8, 2.6, 3.6, 3.8)
add_gt_symbol(slide, 8.9, 2.6, 40, PURPLE)
add_text_box(slide, 8.9, 3.3, 3.2, 0.4, 'What it cannot do', font_size=20, bold=True, color=BLACK)
add_text_box(slide, 8.9, 3.8, 3.2, 1.5, 'It doesn\'t "think" like humans. It predicts responses based on patterns. It can make mistakes and should always be reviewed.', font_size=15, color=MID_GRAY)

# ============================================================
# SLIDE 3: Meet Claude
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, 'Meet Claude', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)

add_text_box(slide, 0.8, 1.4, 8, 0.5, 'Claude is an AI assistant built by Anthropic — designed to be helpful, harmless, and honest.', font_size=18, color=MID_GRAY)

# Left: key attributes
attributes = [
    ('Helpful', 'Designed to give accurate, useful answers'),
    ('Harmless', 'Built with safety guardrails and alignment'),
    ('Honest', 'Tells you when it\'s unsure or doesn\'t know'),
    ('Versatile', 'Writes code, explains concepts, creates content'),
    ('Conversational', 'Remembers context within a session'),
]

for i, (title, desc) in enumerate(attributes):
    y = 2.2 + i * 0.85
    add_rect(slide, 0.8, y, 0.06, 0.55, PURPLE)
    add_text_box(slide, 1.1, y - 0.05, 2.5, 0.4, title, font_size=18, bold=True, color=BLACK)
    add_text_box(slide, 1.1, y + 0.3, 5.5, 0.4, desc, font_size=15, color=MID_GRAY)

# Right: Model family card
add_rect(slide, 7.8, 2.0, 4.8, 4.6, LIGHT_GRAY)
add_text_box(slide, 8.1, 2.2, 4.2, 0.5, 'Claude Model Family', font_size=20, bold=True, color=BLACK)

models = [
    ('Haiku', 'Fast & affordable — great for quick tasks and high-volume use cases', GREEN),
    ('Sonnet', 'Balanced performance — ideal for most development work', BLUE),
    ('Opus', 'Most capable — complex reasoning, architecture, and nuanced tasks', PURPLE),
]

for i, (name, desc, color) in enumerate(models):
    y = 3.0 + i * 1.15
    add_card_white(slide, 8.1, y, 4.2, 0.95)
    add_rect(slide, 8.1, y, 0.08, 0.95, color)
    add_text_box(slide, 8.4, y + 0.05, 2, 0.35, name, font_size=17, bold=True, color=BLACK)
    add_text_box(slide, 8.4, y + 0.4, 3.7, 0.5, desc, font_size=13, color=MID_GRAY)

# ============================================================
# SLIDE 4: Claude vs Claude Code
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, 'Claude vs Claude Code', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)

# Left column
add_rect(slide, 0.8, 1.6, 5.5, 5.0, LIGHT_GRAY)
add_text_box(slide, 1.2, 1.8, 4.8, 0.5, 'Claude', font_size=28, bold=True, color=BLACK)
add_text_box(slide, 1.2, 2.3, 4.8, 0.4, 'The Chatbot', font_size=16, color=PURPLE)

items_claude = [
    '>  Chat on claude.ai or the mobile app',
    '>  Ask questions, get answers',
    '>  Copy & paste code manually',
    '>  Great for learning & research',
    '>  Like texting a smart friend',
]
for i, item in enumerate(items_claude):
    add_text_box(slide, 1.2, 3.0 + i * 0.6, 4.8, 0.5, item, font_size=16, color=DARK_GRAY)

# Right column
add_rect(slide, 7.0, 1.6, 5.5, 5.0, PURPLE)
add_text_box(slide, 7.4, 1.8, 4.8, 0.5, 'Claude Code', font_size=28, bold=True, color=WHITE)
add_text_box(slide, 7.4, 2.3, 4.8, 0.4, 'The Developer Tool', font_size=16, color=WHITE)

items_cc = [
    '>  Works directly in your terminal',
    '>  Reads & edits your files',
    '>  Runs commands for you',
    '>  Builds entire applications',
    '>  Like a developer next to you',
]
for i, item in enumerate(items_cc):
    add_text_box(slide, 7.4, 3.0 + i * 0.6, 4.8, 0.5, item, font_size=16, color=WHITE)

# Arrow between
add_text_box(slide, 5.8, 3.2, 1.5, 1, '>', font_size=60, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: What Can You Build?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8, 'What Can You Build?', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)
add_text_box(slide, 0.8, 1.4, 11, 0.5, 'Even as a complete beginner, you can create real, working projects today.', font_size=18, color=MID_GRAY)

grid = [
    ('Websites', 'Personal portfolios, landing\npages, and blogs'),
    ('Web Apps', 'Dashboards, tools, and\ninteractive applications'),
    ('APIs', 'Backend services that\npower applications'),
    ('Automations', 'Scripts and tools that\nsave hours of work'),
    ('Games', 'Browser-based games\nwith HTML5 & JavaScript'),
    ('AI Apps', 'Chatbots, tutors, and\nAI-powered tools'),
]

colors = [PURPLE, DARK_PURPLE, BLUE, GREEN, RGBColor(0xF5, 0x82, 0x20), RGBColor(0xE0, 0x00, 0x77)]

for i, ((title, desc), color) in enumerate(zip(grid, colors)):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 2.2 + row * 2.4

    add_card_white(slide, x, y, 3.7, 2.0)
    add_rect(slide, x, y, 3.7, 0.06, color)
    add_gt_symbol(slide, x + 0.15, y + 0.15, 28, color)
    add_text_box(slide, x + 0.6, y + 0.25, 2.8, 0.4, title, font_size=20, bold=True, color=BLACK)
    add_text_box(slide, x + 0.3, y + 0.8, 3.1, 0.9, desc, font_size=15, color=MID_GRAY)

# ============================================================
# SLIDE 6: How Claude Code Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8, 'How Claude Code Works', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)
add_text_box(slide, 0.8, 1.4, 11, 0.5, 'A simple four-step workflow: describe, generate, review, refine.', font_size=18, color=MID_GRAY)

steps = [
    ('01', 'Describe', 'Tell Claude Code what\nyou want in plain English'),
    ('02', 'Generate', 'Claude writes the code\nand creates files'),
    ('03', 'Review', 'Check the proposed\nchanges and approve'),
    ('04', 'Refine', 'Ask for adjustments\nuntil it\'s perfect'),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.8 + i * 3.1

    # Number with purple background
    add_rect(slide, x, 2.2, 2.8, 0.7, PURPLE)
    add_text_box(slide, x, 2.25, 0.8, 0.6, num, font_size=24, bold=True, color=WHITE)
    add_text_box(slide, x + 0.7, 2.25, 2, 0.6, title, font_size=22, bold=True, color=WHITE)

    # Description
    add_rect(slide, x, 2.9, 2.8, 1.3, LIGHT_GRAY)
    add_text_box(slide, x + 0.2, 3.1, 2.4, 0.9, desc, font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

    # Arrow between steps
    if i < 3:
        add_text_box(slide, x + 2.7, 2.6, 0.5, 0.5, '>', font_size=28, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

# Example box
add_rect(slide, 0.8, 4.8, 11.7, 1.8, LIGHT_GRAY)
add_rect(slide, 0.8, 4.8, 0.08, 1.8, PURPLE)
add_text_box(slide, 1.2, 4.9, 3, 0.4, 'EXAMPLE', font_size=14, bold=True, color=PURPLE)
add_text_box(slide, 1.2, 5.3, 10.5, 0.4, 'You say:', font_size=15, bold=True, color=BLACK)
add_text_box(slide, 1.2, 5.7, 10.5, 0.4, '"Create a personal portfolio webpage with my name, an about section, and a contact form"', font_size=15, color=DARK_GRAY)
add_text_box(slide, 1.2, 6.1, 10.5, 0.4, 'Claude Code creates the complete HTML, CSS, and JavaScript — ready to open in your browser.', font_size=15, bold=True, color=PURPLE)

# ============================================================
# SLIDE 7: Key Concepts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, 'Key Concepts', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)

concepts = [
    ('Prompts', 'The messages you send to Claude Code. Better prompts = better results.', 'Be specific: "Create an HTML page with a blue header" beats "make a website"'),
    ('Context', 'Claude Code sees the files in your current folder and uses them for better answers.', 'Always start Claude Code in the right project folder'),
    ('Permissions', 'Claude Code asks before changing any of your files. You stay in control.', 'Always review changes before approving them'),
    ('Slash Commands', 'Built-in shortcuts: /help, /clear, /compact to manage your session.', 'Use /compact when conversations get long to free up memory'),
]

for i, (title, desc, tip) in enumerate(concepts):
    y = 1.5 + i * 1.35
    add_rect(slide, 0.8, y, 11.7, 1.15, LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_rect(slide, 0.8, y, 0.08, 1.15, PURPLE)

    add_text_box(slide, 1.2, y + 0.1, 2, 0.35, title, font_size=20, bold=True, color=BLACK)
    add_text_box(slide, 1.2, y + 0.5, 5.5, 0.5, desc, font_size=15, color=DARK_GRAY)

    # Tip on the right
    add_rect(slide, 8.5, y + 0.15, 0.06, 0.8, GREEN)
    add_text_box(slide, 8.8, y + 0.1, 3.5, 0.9, tip, font_size=14, color=MID_GRAY)

# ============================================================
# SLIDE 8: Bootcamp Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_purple_bar(slide)

add_text_box(slide, 0.8, 0.4, 11, 0.8, 'Today\'s Bootcamp Agenda', font_size=36, bold=True, color=BLACK)
add_rect(slide, 0.8, 1.1, 1.2, 0.05, PURPLE)
add_text_box(slide, 0.8, 1.4, 11, 0.5, '2.5 hours of hands-on learning with an AI tutor guiding you every step.', font_size=18, color=MID_GRAY)

modules = [
    ('01', 'What is AI & Claude Code?', '15 min'),
    ('02', 'Setting Up Claude Code', '15 min'),
    ('03', 'Your First Conversation', '15 min'),
    ('04', 'Prompts & Skills', '20 min'),
    ('05', 'File Operations', '20 min'),
    ('06', 'Building a Web Page', '25 min'),
    ('07', 'Git & GitHub', '20 min'),
    ('08', 'Tips & Next Steps', '20 min'),
]

for i, (num, title, duration) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.1 + row * 1.2

    bg_color = LIGHT_GRAY if (i % 2 == 0) else WHITE
    add_rect(slide, x, y, 5.8, 0.95, bg_color)
    add_rect(slide, x, y, 0.06, 0.95, PURPLE)

    add_text_box(slide, x + 0.3, y + 0.2, 0.6, 0.5, num, font_size=22, bold=True, color=PURPLE)
    add_text_box(slide, x + 1.0, y + 0.25, 3.5, 0.5, title, font_size=17, color=BLACK)
    add_text_box(slide, x + 4.6, y + 0.28, 1.0, 0.4, duration, font_size=14, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

# Total time
add_rect(slide, 0.8, 6.2, 11.7, 0.5, PURPLE)
add_text_box(slide, 1.2, 6.22, 5, 0.4, 'Total Duration', font_size=16, bold=True, color=WHITE)
add_text_box(slide, 9.5, 6.22, 2.8, 0.4, '2 hours 30 minutes', font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 9: Let's Get Started
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PURPLE)

# Large > background
add_text_box(slide, 8.5, -1, 5, 10, '>', font_size=400, bold=True, color=DARK_PURPLE, alignment=PP_ALIGN.RIGHT)

add_text_box(slide, 1, 1.5, 8, 1, 'Let\'s get started.', font_size=48, bold=True, color=WHITE)

add_rect(slide, 1, 2.8, 2, 0.04, WHITE)

add_text_box(slide, 1, 3.2, 8, 0.5, 'Open your browser and navigate to:', font_size=20, color=WHITE)

# URL box
add_rect(slide, 1, 4.0, 6, 0.9, WHITE)
add_text_box(slide, 1.3, 4.1, 5.5, 0.7, 'emcidie.com/bootcamp/', font_size=32, bold=True, color=PURPLE)

add_text_box(slide, 1, 5.2, 8, 0.5, '>  Enter your name to begin', font_size=20, color=WHITE)
add_text_box(slide, 1, 5.7, 8, 0.5, '>  Your AI tutor will guide you through each module', font_size=20, color=WHITE)
add_text_box(slide, 1, 6.2, 8, 0.5, '>  Ask questions anytime — no question is too basic', font_size=20, color=WHITE)

# Save
output_path = '/Users/mikeduff/claude-code-bootcamp/Japan_CMT_Introduction_to_Claude_Accenture.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
